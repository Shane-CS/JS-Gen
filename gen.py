# Required Libraries
import os
from pptx import Presentation
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from sklearn.feature_extraction.text import TfidfVectorizer
from transformers import GPT2LMHeadModel, GPT2Tokenizer

# Step 1: Load the presentations
ppts = []
for file in os.listdir('/path/to/your/ppts'):
    if file.endswith(".ppt") or file.endswith(".pptx"):
        ppts.append(Presentation(os.path.join('/path/to/your/ppts', file)))

# Step 2: Extract text
texts = []
for ppt in ppts:
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        texts.append(run.text)

# Step 3: Content Breakdown
tokens = []
for text in texts:
    tokens.extend(word_tokenize(text))  # Tokenization
tokens = [word for word in tokens if word.isalnum()]  # Remove punctuation
tokens = [word.lower() for word in tokens if not word.lower() in stopwords.words()]  # Remove stopwords
stemmer = PorterStemmer()
tokens = [stemmer.stem(token) for token in tokens]  # Apply stemming

# Step 4: Course Breakdown
# 'courses' is a list of documents with each representing a course
courses = []
for file in os.listdir('/path/to/your/courses'):
    with open(file, 'r') as course_file:
        courses.append(course_file.read())

# Step 5: Content Matching
vectorizer = TfidfVectorizer()
vectorizer.fit_transform(courses + [' '.join(tokens)])  # Update the vocab
courses_tfidf = vectorizer.transform(courses)
tokens_tfidf = vectorizer.transform([' '.join(tokens)])

# Step 6: Script Generation
tokenizer = GPT2Tokenizer.from_pretrained("gpt2")
model = GPT2LMHeadModel.from_pretrained("gpt2")

# for each course, generate the script using GPT-2 or GPT-3
for i, course in enumerate(courses):
    inputs = tokenizer.encode(course, return_tensors='pt')
    outputs = model.generate(inputs, max_length=500, temperature=0.7, num_return_sequences=1)
    print("Course", i, "Script:", tokenizer.decode(outputs[0]))